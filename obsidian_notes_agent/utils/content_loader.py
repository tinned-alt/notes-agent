"""Content loading utilities for various file formats and URLs."""

import re
from pathlib import Path
from typing import Dict, Optional, List
from urllib.parse import urlparse, parse_qs
import requests
from bs4 import BeautifulSoup


class ContentLoader:
    """Load and extract content from various sources."""

    # Standard Chrome User-Agent for web requests
    USER_AGENT = (
        'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 '
        '(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
    )
    # HTTP request timeout in seconds
    REQUEST_TIMEOUT = 15

    @staticmethod
    def _format_timestamp(total_seconds: int) -> str:
        """Format seconds into timestamp string.

        Uses HH:MM:SS for videos >= 1 hour, MM:SS otherwise.

        Args:
            total_seconds: Number of seconds

        Returns:
            Formatted timestamp string
        """
        hours = total_seconds // 3600
        minutes = (total_seconds % 3600) // 60
        seconds = total_seconds % 60

        if hours > 0:
            return f"[{hours:d}:{minutes:02d}:{seconds:02d}]"
        return f"[{minutes:02d}:{seconds:02d}]"

    @staticmethod
    def _format_duration(total_seconds: int) -> str:
        """Format total seconds into duration string.

        Uses H:MM:SS for durations >= 1 hour, M:SS otherwise.

        Args:
            total_seconds: Total duration in seconds

        Returns:
            Formatted duration string
        """
        hours = total_seconds // 3600
        minutes = (total_seconds % 3600) // 60
        seconds = total_seconds % 60

        if hours > 0:
            return f"{hours}:{minutes:02d}:{seconds:02d}"
        return f"{minutes}:{seconds:02d}"

    @staticmethod
    def extract_youtube_video_id(url: str) -> Optional[str]:
        """Extract YouTube video ID from various URL formats.

        Supports:
            - https://www.youtube.com/watch?v=VIDEO_ID
            - https://youtu.be/VIDEO_ID
            - https://www.youtube.com/embed/VIDEO_ID
            - https://www.youtube.com/v/VIDEO_ID
            - https://m.youtube.com/watch?v=VIDEO_ID

        Args:
            url: YouTube URL

        Returns:
            Video ID string or None if not found
        """
        # Pattern for youtu.be short URLs
        short_pattern = r'(?:youtu\.be/)([a-zA-Z0-9_-]{11})'
        # Pattern for standard youtube.com URLs
        standard_pattern = r'(?:youtube\.com/(?:watch\?v=|embed/|v/))([a-zA-Z0-9_-]{11})'

        # Try short URL pattern first
        match = re.search(short_pattern, url)
        if match:
            return match.group(1)

        # Try standard URL pattern
        match = re.search(standard_pattern, url)
        if match:
            return match.group(1)

        # Try parsing query string for watch URLs
        parsed = urlparse(url)
        if 'youtube.com' in parsed.netloc:
            query_params = parse_qs(parsed.query)
            if 'v' in query_params:
                return query_params['v'][0]

        return None

    @staticmethod
    def is_youtube_url(url: str) -> bool:
        """Check if a URL is a YouTube video URL.

        Args:
            url: URL to check

        Returns:
            True if URL is a YouTube video URL
        """
        parsed = urlparse(url)
        youtube_domains = ['youtube.com', 'www.youtube.com', 'm.youtube.com', 'youtu.be']
        return parsed.netloc in youtube_domains

    @staticmethod
    def load_from_youtube(url: str) -> Dict[str, any]:
        """Load transcript from a YouTube video.

        Args:
            url: YouTube video URL

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
            from youtube_transcript_api._errors import (
                TranscriptsDisabled,
                NoTranscriptFound,
                VideoUnavailable
            )
        except ImportError:
            raise ImportError(
                "youtube-transcript-api is required for YouTube processing. "
                "Install with: pip install youtube-transcript-api"
            )

        video_id = ContentLoader.extract_youtube_video_id(url)
        if not video_id:
            raise ValueError(f"Could not extract video ID from URL: {url}")

        try:
            # Fetch video title from YouTube page
            headers = {'User-Agent': ContentLoader.USER_AGENT}
            response = requests.get(url, headers=headers, timeout=ContentLoader.REQUEST_TIMEOUT)
            soup = BeautifulSoup(response.text, 'html.parser')

            # Extract title
            title_tag = soup.find('title')
            title = title_tag.get_text().strip() if title_tag else f"YouTube Video {video_id}"
            # Remove " - YouTube" suffix if present
            if title.endswith(' - YouTube'):
                title = title[:-10]

            # Extract channel name if available
            channel = None
            channel_tag = soup.find('link', {'itemprop': 'name'})
            if channel_tag and channel_tag.get('content'):
                channel = channel_tag.get('content')

        except Exception:
            title = f"YouTube Video {video_id}"
            channel = None

        try:
            # Get transcript - try to get English first, then fall back to any available
            transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)

            transcript = None
            # Try English first
            try:
                transcript = transcript_list.find_transcript(['en', 'en-US', 'en-GB'])
            except NoTranscriptFound:
                # Fall back to any available transcript (auto-generated or manual)
                try:
                    transcript = transcript_list.find_generated_transcript(['en'])
                except NoTranscriptFound:
                    # Get any available transcript
                    for t in transcript_list:
                        transcript = t
                        break

            if transcript is None:
                raise NoTranscriptFound(video_id, ['en'], None)

            # Fetch the actual transcript data
            transcript_data = transcript.fetch()

            # Format transcript with timestamps (HH:MM:SS for videos >= 1 hour)
            content_parts = []
            for entry in transcript_data:
                seconds = int(entry['start'])
                timestamp = ContentLoader._format_timestamp(seconds)
                text = entry['text'].strip()
                content_parts.append(f"{timestamp} {text}")

            content = "\n".join(content_parts)

            # Calculate duration
            if transcript_data:
                last_entry = transcript_data[-1]
                total_seconds = int(last_entry['start'] + last_entry.get('duration', 0))
                duration = ContentLoader._format_duration(total_seconds)
            else:
                duration = "Unknown"

            metadata = {
                'url': url,
                'video_id': video_id,
                'channel': channel,
                'duration': duration,
                'transcript_language': transcript.language,
                'is_generated': transcript.is_generated,
            }

            return {
                'title': title,
                'content': content,
                'source': url,
                'source_type': 'youtube',
                'metadata': metadata
            }

        except TranscriptsDisabled:
            raise ValueError(f"Transcripts are disabled for video: {url}")
        except VideoUnavailable:
            raise ValueError(f"Video is unavailable: {url}")
        except NoTranscriptFound:
            raise ValueError(f"No transcript available for video: {url}")
        except Exception as e:
            raise ValueError(f"Failed to load YouTube transcript from {url}: {str(e)}")

    @staticmethod
    def load_from_url(url: str) -> Dict[str, any]:
        """Load content from a URL.

        Args:
            url: URL to fetch content from

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        try:
            # Fetch the URL
            headers = {'User-Agent': ContentLoader.USER_AGENT}
            response = requests.get(url, headers=headers, timeout=ContentLoader.REQUEST_TIMEOUT)
            response.raise_for_status()

            # Parse HTML
            soup = BeautifulSoup(response.text, 'html.parser')

            # Extract title
            title = soup.find('title')
            title = title.get_text().strip() if title else urlparse(url).path.strip('/')

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            # Get text content
            text = soup.get_text()

            # Clean up text
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            content = '\n'.join(chunk for chunk in chunks if chunk)

            return {
                'title': title,
                'content': content,
                'source': url,
                'source_type': 'url',
                'metadata': {
                    'url': url,
                    'domain': urlparse(url).netloc
                }
            }

        except Exception as e:
            raise ValueError(f"Failed to load URL {url}: {str(e)}")

    @staticmethod
    def load_from_pdf(file_path: str) -> Dict[str, any]:
        """Load content from a PDF file.

        Args:
            file_path: Path to PDF file

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        try:
            from pypdf import PdfReader

            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            reader = PdfReader(str(path))

            # Extract text from all pages
            content = ""
            for page in reader.pages:
                content += page.extract_text() + "\n\n"

            # Get metadata
            metadata = {}
            if reader.metadata:
                metadata = {
                    'author': reader.metadata.get('/Author', ''),
                    'title': reader.metadata.get('/Title', ''),
                    'subject': reader.metadata.get('/Subject', ''),
                    'creator': reader.metadata.get('/Creator', ''),
                }

            title = metadata.get('title') or path.stem

            return {
                'title': title,
                'content': content.strip(),
                'source': str(path),
                'source_type': 'pdf',
                'metadata': metadata
            }

        except ImportError:
            raise ImportError("pypdf is required for PDF processing. Install with: pip install pypdf")
        except Exception as e:
            raise ValueError(f"Failed to load PDF {file_path}: {str(e)}")

    @staticmethod
    def load_from_docx(file_path: str) -> Dict[str, any]:
        """Load content from a DOCX file.

        Args:
            file_path: Path to DOCX file

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        try:
            from docx import Document

            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            doc = Document(str(path))

            # Extract all paragraphs
            content = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])

            # Extract metadata
            core_properties = doc.core_properties
            metadata = {
                'author': core_properties.author or '',
                'title': core_properties.title or '',
                'subject': core_properties.subject or '',
                'keywords': core_properties.keywords or '',
            }

            title = metadata.get('title') or path.stem

            return {
                'title': title,
                'content': content,
                'source': str(path),
                'source_type': 'docx',
                'metadata': metadata
            }

        except ImportError:
            raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
        except Exception as e:
            raise ValueError(f"Failed to load DOCX {file_path}: {str(e)}")

    @staticmethod
    def load_from_pptx(file_path: str) -> Dict[str, any]:
        """Load content from a PowerPoint file.

        Args:
            file_path: Path to PPTX file

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        try:
            from pptx import Presentation

            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            prs = Presentation(str(path))

            # Extract text from all slides
            content_parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"## Slide {i}\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                if slide_text.strip() != f"## Slide {i}":
                    content_parts.append(slide_text)

            content = "\n\n".join(content_parts)

            # Get metadata
            metadata = {
                'author': prs.core_properties.author or '',
                'title': prs.core_properties.title or '',
                'subject': prs.core_properties.subject or '',
            }

            title = metadata.get('title') or path.stem

            return {
                'title': title,
                'content': content,
                'source': str(path),
                'source_type': 'pptx',
                'metadata': metadata
            }

        except ImportError:
            raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
        except Exception as e:
            raise ValueError(f"Failed to load PPTX {file_path}: {str(e)}")

    @staticmethod
    def load_from_markdown(file_path: str) -> Dict[str, any]:
        """Load content from a Markdown file.

        Args:
            file_path: Path to Markdown file

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Try to extract title from first heading
        lines = content.split('\n')
        title = path.stem
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
                break

        return {
            'title': title,
            'content': content,
            'source': str(path),
            'source_type': 'markdown',
            'metadata': {}
        }

    @staticmethod
    def load_from_text(file_path: str) -> Dict[str, any]:
        """Load content from a plain text file.

        Args:
            file_path: Path to text file

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        # Try to extract title from first non-empty line
        lines = content.split('\n')
        title = path.stem
        for line in lines:
            if line.strip():
                title = line.strip()[:100]  # Use first 100 chars of first line
                break

        return {
            'title': title,
            'content': content,
            'source': str(path),
            'source_type': 'text',
            'metadata': {}
        }

    @staticmethod
    def load_from_rtf(file_path: str) -> Dict[str, any]:
        """Load content from an RTF file.

        Args:
            file_path: Path to RTF file

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        try:
            from striprtf.striprtf import rtf_to_text

            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")

            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()

            content = rtf_to_text(rtf_content)

            # Try to extract title from first line
            lines = content.split('\n')
            title = path.stem
            for line in lines:
                if line.strip():
                    title = line.strip()[:100]
                    break

            return {
                'title': title,
                'content': content,
                'source': str(path),
                'source_type': 'rtf',
                'metadata': {}
            }

        except ImportError:
            raise ImportError("striprtf is required for RTF processing. Install with: pip install striprtf")
        except Exception as e:
            raise ValueError(f"Failed to load RTF {file_path}: {str(e)}")

    @staticmethod
    def load_content(source: str) -> Dict[str, any]:
        """Load content from a URL or file path.

        Automatically detects the source type and uses the appropriate loader.

        Args:
            source: URL or file path

        Returns:
            Dictionary with 'title', 'content', 'source', and 'metadata'
        """
        # Check if it's a URL
        if source.startswith(('http://', 'https://')):
            # Check if it's a YouTube URL
            if ContentLoader.is_youtube_url(source):
                return ContentLoader.load_from_youtube(source)
            return ContentLoader.load_from_url(source)

        # Otherwise treat as file path
        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"File or URL not found: {source}")

        # Determine file type by extension
        extension = path.suffix.lower()

        loaders = {
            '.pdf': ContentLoader.load_from_pdf,
            '.docx': ContentLoader.load_from_docx,
            '.doc': ContentLoader.load_from_docx,
            '.pptx': ContentLoader.load_from_pptx,
            '.ppt': ContentLoader.load_from_pptx,
            '.md': ContentLoader.load_from_markdown,
            '.markdown': ContentLoader.load_from_markdown,
            '.txt': ContentLoader.load_from_text,
            '.rtf': ContentLoader.load_from_rtf,
        }

        loader = loaders.get(extension)
        if not loader:
            # Default to text loader for unknown types
            return ContentLoader.load_from_text(source)

        return loader(source)
